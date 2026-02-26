"""
CropSense AI — TEXPO 2026 Pitch Deck Generator
Generates a professional, dark-themed investor pitch deck (.pptx)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Color Palette (Dark Glassmorphic) ──────────────────────────────────────────
BG_DARK       = RGBColor(0x0B, 0x0F, 0x19)   # Deep navy
BG_CARD       = RGBColor(0x11, 0x18, 0x27)   # Card background
ACCENT_GREEN  = RGBColor(0x34, 0xD3, 0x99)   # Emerald accent
ACCENT_TEAL   = RGBColor(0x14, 0xB8, 0xA6)   # Teal accent
TEXT_WHITE     = RGBColor(0xF1, 0xF5, 0xF9)   # Primary text
TEXT_MUTED     = RGBColor(0x94, 0xA3, 0xB8)   # Muted text
HIGHLIGHT_RED  = RGBColor(0xEF, 0x44, 0x44)   # Warning red
HIGHLIGHT_AMBER= RGBColor(0xF5, 0x9E, 0x0B)   # Amber
HIGHLIGHT_BLUE = RGBColor(0x38, 0xBD, 0xF8)   # Sky blue

SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT

# ── Helper Functions ───────────────────────────────────────────────────────────

def add_bg(slide, color=BG_DARK):
    """Fill entire slide background with a solid dark color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_rect(slide, left, top, width, height, fill_color=BG_CARD, border_color=None, opacity=None):
    """Add a rounded rectangle card."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=TEXT_WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    """Add a text box with styled text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_frame(slide, left, top, width, height, bullets, font_size=16, color=TEXT_WHITE, bullet_color=ACCENT_GREEN):
    """Add a text frame with multiple bullet points."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.space_after = Pt(10)
        p.level = 0
    return txBox

def add_accent_line(slide, left, top, width):
    """Add a thin emerald accent line."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_GREEN
    shape.line.fill.background()
    return shape

def add_stat_card(slide, left, top, number, label, accent=ACCENT_GREEN):
    """Add a stat card with big number and label."""
    card = add_shape_rect(slide, left, top, Inches(2.8), Inches(1.8), BG_CARD, border_color=RGBColor(0x1E, 0x29, 0x3B))
    add_text_box(slide, left + Inches(0.3), top + Inches(0.25), Inches(2.2), Inches(0.9),
                 number, font_size=36, color=accent, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.3), top + Inches(1.1), Inches(2.2), Inches(0.5),
                 label, font_size=13, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

def add_icon_bullet(slide, left, top, icon_text, desc, icon_color=ACCENT_GREEN):
    """Add an icon-style bullet: [icon] Description."""
    add_text_box(slide, left, top, Inches(0.5), Inches(0.4), icon_text, font_size=20, color=icon_color, bold=True)
    add_text_box(slide, left + Inches(0.5), top, Inches(5), Inches(0.4), desc, font_size=15, color=TEXT_WHITE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE & HOOK
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
add_bg(slide)

# Decorative gradient bar at top
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08))
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT_GREEN
bar.line.fill.background()

# Title
add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
             "🌿 CropSense AI", font_size=54, color=ACCENT_GREEN, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(3.2), Inches(11), Inches(0.8),
             "Hyper-Local Crop Disease Detection & Advisory Platform", font_size=28, color=TEXT_WHITE, alignment=PP_ALIGN.CENTER)

add_accent_line(slide, Inches(5), Inches(4.2), Inches(3.3))

add_text_box(slide, Inches(1), Inches(4.6), Inches(11), Inches(0.6),
             "TEXPO 2026  •  Investor Pitch", font_size=20, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.5),
             "AI + Computer Vision + Real-Time Weather Intelligence", font_size=16, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.7),
             "The Problem", font_size=36, color=HIGHLIGHT_RED, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(2))

# Stat cards row
add_stat_card(slide, Inches(0.8), Inches(1.8), "$220B+", "Annual Global\nCrop Losses", HIGHLIGHT_RED)
add_stat_card(slide, Inches(4.0), Inches(1.8), "80%", "Farmers Misdiagnose\nPlant Diseases", HIGHLIGHT_AMBER)
add_stat_card(slide, Inches(7.2), Inches(1.8), "0%", "Weather-Aware\nDiagnosis Tools", HIGHLIGHT_BLUE)

# Problem bullets card
card = add_shape_rect(slide, Inches(0.8), Inches(4.1), Inches(11.5), Inches(3),
                      BG_CARD, border_color=RGBColor(0x1E, 0x29, 0x3B))
bullets = [
    "⚠  Farmers rely on visual guesswork → delayed treatment → entire harvests lost",
    "⚠  No existing tool connects disease identification with local weather spread risk",
    "⚠  Rural farmers lack access to English-only, complex diagnostic platforms",
    "⚠  Government agencies have zero real-time outbreak tracking dashboards"
]
add_bullet_frame(slide, Inches(1.2), Inches(4.5), Inches(10.5), Inches(2.5), bullets, font_size=16, color=TEXT_WHITE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — THE SOLUTION
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.7),
             "Our Solution: CropSense AI", font_size=36, color=ACCENT_GREEN, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(3))

# Solution pillars
pillars = [
    ("🔬", "Instant Disease Diagnosis", "Upload a leaf photo → AI identifies the disease in seconds with 95%+ accuracy"),
    ("🌦️", "48-Hour Spread Risk", "Live weather data calculates how fast the disease will spread in your exact location"),
    ("🗺️", "Outbreak Heatmap", "Every detection is plotted on a real-time map — a dashboard for agricultural agencies"),
    ("🧠", "Explainable AI (Grad-CAM)", "Shows farmers exactly WHICH spots on the leaf triggered the diagnosis"),
    ("🌐", "Multilingual Access", "Full support for English, Hindi, and Telugu — built for rural deployment"),
    ("💊", "Treatment Advisories", "Instant, actionable remedies delivered in the farmer's native language"),
]

for i, (icon, title, desc) in enumerate(pillars):
    y = Inches(1.7) + Inches(i * 0.9)
    add_text_box(slide, Inches(1), y, Inches(0.5), Inches(0.5), icon, font_size=22, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.6), y, Inches(3), Inches(0.4), title, font_size=18, color=ACCENT_GREEN, bold=True)
    add_text_box(slide, Inches(1.6), y + Inches(0.35), Inches(10), Inches(0.4), desc, font_size=14, color=TEXT_MUTED)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — PRODUCT DEMO / USER JOURNEY
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.7),
             "Product Demo: User Journey", font_size=36, color=ACCENT_GREEN, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(3))

steps = [
    ("Step 1", "Login / Register", "Secure JWT-based authentication with multi-language UI"),
    ("Step 2", "Upload Leaf Photo", "Drag & drop or camera capture — mobile optimized"),
    ("Step 3", "AI Processing", "OpenCV DIP pipeline preprocesses → MobileNetV3 CNN classifies the disease"),
    ("Step 4", "Results Dashboard", "Disease name, confidence %, severity, treatment, and Grad-CAM overlay"),
    ("Step 5", "Spread Risk Alert", "Live geolocation + OpenWeatherMap → 48hr spread probability"),
    ("Step 6", "Outbreak Map", "Detection logged on the real-time Leaflet heatmap for all users"),
]

for i, (step, title, desc) in enumerate(steps):
    col = i % 3
    row = i // 3
    x = Inches(0.8) + Inches(col * 4.1)
    y = Inches(1.7) + Inches(row * 2.6)
    card = add_shape_rect(slide, x, y, Inches(3.7), Inches(2.2), BG_CARD, border_color=RGBColor(0x1E, 0x29, 0x3B))
    add_text_box(slide, x + Inches(0.2), y + Inches(0.15), Inches(1), Inches(0.3),
                 step, font_size=11, color=ACCENT_TEAL, bold=True)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.5), Inches(3.2), Inches(0.4),
                 title, font_size=18, color=TEXT_WHITE, bold=True)
    add_text_box(slide, x + Inches(0.2), y + Inches(1.05), Inches(3.2), Inches(1),
                 desc, font_size=13, color=TEXT_MUTED)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — SECRET SAUCE #1: DIP PIPELINE
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
             "Secret Sauce #1: Digital Image Processing Pipeline", font_size=32, color=ACCENT_GREEN, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(4))

add_text_box(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(0.5),
             "\"Pure deep learning is noisy. We built a custom OpenCV preprocessing layer that boosted accuracy by +28%.\"",
             font_size=16, color=HIGHLIGHT_AMBER, bold=False)

# Pipeline steps
pipeline = [
    ("1️⃣", "Median Filtering", "Removes sensor noise and artifacts from camera-captured leaf images"),
    ("2️⃣", "CLAHE Enhancement", "Contrast Limited Adaptive Histogram Equalization — amplifies leaf veins and disease spots in LAB color space"),
    ("3️⃣", "HSV Segmentation", "Forcibly suppresses background (soil, sky, hands) — feeds the CNN ONLY the relevant botanical tissue"),
    ("4️⃣", "MobileNetV3 CNN", "The cleaned, enhanced image enters the deep learning classifier for disease prediction"),
]

for i, (num, title, desc) in enumerate(pipeline):
    y = Inches(2.5) + Inches(i * 1.15)
    add_shape_rect(slide, Inches(0.8), y, Inches(11.5), Inches(0.95),
                   BG_CARD, border_color=RGBColor(0x1E, 0x29, 0x3B))
    add_text_box(slide, Inches(1.1), y + Inches(0.1), Inches(0.5), Inches(0.4), num, font_size=18, color=ACCENT_GREEN, bold=True)
    add_text_box(slide, Inches(1.7), y + Inches(0.08), Inches(3), Inches(0.4), title, font_size=17, color=TEXT_WHITE, bold=True)
    add_text_box(slide, Inches(1.7), y + Inches(0.48), Inches(10), Inches(0.4), desc, font_size=13, color=TEXT_MUTED)

# Result highlight
add_stat_card(slide, Inches(10.2), Inches(1.4), "+28%", "Accuracy Boost\nvs Raw Model", ACCENT_GREEN)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — SECRET SAUCE #2: SPREAD RISK & HEATMAP
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
             "Secret Sauce #2: Predictive Spread Risk & Live Heatmap", font_size=32, color=ACCENT_GREEN, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(4))

# Left column - Spread Risk
add_shape_rect(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5),
               BG_CARD, border_color=RGBColor(0x1E, 0x29, 0x3B))
add_text_box(slide, Inches(1.2), Inches(2.0), Inches(4.8), Inches(0.5),
             "🌦️ 48-Hour Spread Risk Engine", font_size=20, color=HIGHLIGHT_AMBER, bold=True)
bullets_left = [
    "✦ Auto-detects user geolocation via browser GPS",
    "✦ Fetches real-time temperature & humidity from OpenWeatherMap API",
    "✦ Calculates dynamic spread probability using disease-specific epidemiological factors",
    "✦ Alerts farmer: LOW / MEDIUM / HIGH risk with actionable advice",
    "✦ Covers the next 48-hour forecast window"
]
add_bullet_frame(slide, Inches(1.2), Inches(2.7), Inches(4.8), Inches(3.5), bullets_left, font_size=14, color=TEXT_WHITE)

# Right column - Heatmap
add_shape_rect(slide, Inches(6.8), Inches(1.8), Inches(5.5), Inches(5),
               BG_CARD, border_color=RGBColor(0x1E, 0x29, 0x3B))
add_text_box(slide, Inches(7.2), Inches(2.0), Inches(4.8), Inches(0.5),
             "🗺️ Real-Time Outbreak Heatmap", font_size=20, color=HIGHLIGHT_BLUE, bold=True)
bullets_right = [
    "✦ Every detection is georeferenced and stored in database",
    "✦ Interactive Leaflet.js map with clustering",
    "✦ Click any marker → disease name, severity, date",
    "✦ Simulates government agricultural dashboards",
    "✦ Enables macro-level outbreak monitoring at scale"
]
add_bullet_frame(slide, Inches(7.2), Inches(2.7), Inches(4.8), Inches(3.5), bullets_right, font_size=14, color=TEXT_WHITE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — EXPLAINABLE AI: GRAD-CAM
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
             "Explainable AI: Building Trust with Farmers", font_size=36, color=ACCENT_GREEN, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(3))

add_text_box(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(0.8),
             "We don't use a \"black box.\" Grad-CAM shows farmers exactly WHICH spots\non the leaf triggered the AI's neural pathways.",
             font_size=18, color=TEXT_MUTED)

# Grad-CAM explanation cards
cards = [
    ("How It Works", [
        "✦ Extracts activation maps from the final convolutional layer (conv_1)",
        "✦ Weights each activation by gradient importance for the predicted class",
        "✦ Generates a thermal heatmap overlayed on the original leaf image",
        "✦ Red/Yellow zones = highest disease signal regions"
    ]),
    ("Why It Matters", [
        "✦ Farmers can visually verify the AI is looking at the correct spots",
        "✦ Builds trust in AI-driven recommendations in rural communities",
        "✦ Enables agricultural researchers to validate model behavior",
        "✦ Differentiates CropSense AI from generic classification apps"
    ]),
]

for i, (title, bullets) in enumerate(cards):
    x = Inches(0.8) + Inches(i * 6.2)
    add_shape_rect(slide, x, Inches(2.8), Inches(5.8), Inches(4), BG_CARD, border_color=RGBColor(0x1E, 0x29, 0x3B))
    add_text_box(slide, x + Inches(0.3), Inches(3.0), Inches(5), Inches(0.5),
                 title, font_size=20, color=ACCENT_GREEN, bold=True)
    add_bullet_frame(slide, x + Inches(0.3), Inches(3.6), Inches(5.2), Inches(3), bullets, font_size=14, color=TEXT_WHITE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — TECH ARCHITECTURE
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
             "Technical Architecture", font_size=36, color=ACCENT_GREEN, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(2.5))

# Three architecture columns
arch_cols = [
    ("Frontend", HIGHLIGHT_BLUE, [
        "Next.js 15 (App Router)",
        "React + TypeScript",
        "Tailwind CSS + Framer Motion",
        "React-Leaflet (Heatmap)",
        "Axios (API Client)",
        "Responsive & Mobile-First"
    ]),
    ("Backend", ACCENT_GREEN, [
        "FastAPI (Async Python)",
        "SQLAlchemy ORM",
        "PostgreSQL / SQLite",
        "PyJWT Authentication",
        "OpenWeatherMap API",
        "Sarvam AI (Hindi/Telugu TTS)"
    ]),
    ("AI / ML Pipeline", HIGHLIGHT_AMBER, [
        "TensorFlow / Keras",
        "MobileNetV3 (Transfer Learning)",
        "OpenCV DIP Pipeline",
        "Grad-CAM Explainability",
        "TFLite Edge Quantization",
        "Scikit-learn Evaluation"
    ]),
]

for i, (title, color, items) in enumerate(arch_cols):
    x = Inches(0.8) + Inches(i * 4.1)
    card = add_shape_rect(slide, x, Inches(1.7), Inches(3.7), Inches(5.2),
                          BG_CARD, border_color=RGBColor(0x1E, 0x29, 0x3B))
    add_text_box(slide, x + Inches(0.3), Inches(1.9), Inches(3), Inches(0.5),
                 title, font_size=22, color=color, bold=True)
    # Divider
    div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.3), Inches(2.45), Inches(2), Pt(2))
    div.fill.solid()
    div.fill.fore_color.rgb = color
    div.line.fill.background()
    
    for j, item in enumerate(items):
        add_text_box(slide, x + Inches(0.3), Inches(2.7) + Inches(j * 0.6), Inches(3.2), Inches(0.4),
                     f"▸ {item}", font_size=14, color=TEXT_WHITE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — MARKET & IMPACT
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
             "Target Market & Impact Potential", font_size=36, color=ACCENT_GREEN, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(3))

# Stats row
add_stat_card(slide, Inches(0.8), Inches(1.7), "150M+", "Farmers in India\nAlone", ACCENT_GREEN)
add_stat_card(slide, Inches(4.0), Inches(1.7), "$11.2B", "AgriTech Market\nby 2027", HIGHLIGHT_AMBER)
add_stat_card(slide, Inches(7.2), Inches(1.7), "38", "Disease Classes\nSupported", HIGHLIGHT_BLUE)
add_stat_card(slide, Inches(10.4), Inches(1.7), "4.5 MB", "Edge Model\nFootprint", ACCENT_TEAL)

# Market segments
add_shape_rect(slide, Inches(0.8), Inches(4.0), Inches(5.5), Inches(3.2),
               BG_CARD, border_color=RGBColor(0x1E, 0x29, 0x3B))
add_text_box(slide, Inches(1.2), Inches(4.2), Inches(4.8), Inches(0.4),
             "🎯 Target Segments", font_size=18, color=ACCENT_GREEN, bold=True)
segments = [
    "✦ Smallholder farmers in India, SE Asia, and Africa",
    "✦ Government agricultural departments / ICAR",
    "✦ Agricultural cooperatives & NGOs",
    "✦ AgriTech companies seeking disease detection APIs",
]
add_bullet_frame(slide, Inches(1.2), Inches(4.8), Inches(4.8), Inches(2.2), segments, font_size=14, color=TEXT_WHITE)

add_shape_rect(slide, Inches(6.8), Inches(4.0), Inches(5.5), Inches(3.2),
               BG_CARD, border_color=RGBColor(0x1E, 0x29, 0x3B))
add_text_box(slide, Inches(7.2), Inches(4.2), Inches(4.8), Inches(0.4),
             "🌍 Social Impact", font_size=18, color=ACCENT_GREEN, bold=True)
impact = [
    "✦ Reduce crop losses by enabling early detection",
    "✦ Break language barriers for rural communities",
    "✦ Enable data-driven policy for outbreak containment",
    "✦ Edge deployment works in low-connectivity regions",
]
add_bullet_frame(slide, Inches(7.2), Inches(4.8), Inches(4.8), Inches(2.2), impact, font_size=14, color=TEXT_WHITE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — THE ASK & CTA
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1),
             "Thank You", font_size=52, color=ACCENT_GREEN, bold=True, alignment=PP_ALIGN.CENTER)

add_accent_line(slide, Inches(5), Inches(2.7), Inches(3.3))

add_text_box(slide, Inches(1), Inches(3.2), Inches(11), Inches(0.6),
             "CropSense AI — Saving Harvests with Intelligence", font_size=24, color=TEXT_WHITE, alignment=PP_ALIGN.CENTER)

# Contact / key takeaways card
add_shape_rect(slide, Inches(3), Inches(4.2), Inches(7.3), Inches(2.5),
               BG_CARD, border_color=RGBColor(0x1E, 0x29, 0x3B))

takeaways = [
    "✦ Deep Learning + Computer Vision + Real-Time Weather = Smarter Farming",
    "✦ +28% accuracy boost via custom DIP preprocessing pipeline",
    "✦ Edge-ready (4.5 MB) for deployment in low-resource environments",
    "✦ Multilingual, explainable AI designed for real farmers"
]
add_bullet_frame(slide, Inches(3.4), Inches(4.5), Inches(6.5), Inches(2), takeaways, font_size=15, color=TEXT_WHITE)

add_text_box(slide, Inches(1), Inches(7.0), Inches(11), Inches(0.4),
             "TEXPO 2026  •  CropSense AI  •  🌿", font_size=14, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════════════════════
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "CropSenseAI_TEXPO_2026.pptx")
prs.save(output_path)
print(f"\n✅ Presentation saved to: {output_path}")
print(f"📊 Total slides: {len(prs.slides)}")
